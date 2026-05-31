#!/usr/bin/env python3
# 살리자 발표덱 빌더 — 살리자_발표덱_v1.md → 살리자_발표덱.pptx
# 사용법: python3 build_deck.py
# md의 "## S00. 제목" 단위로 슬라이드 생성. 불릿(-), 리드(**굵게**), 인용(>), 코드블럭(```), 발표멘트(> 🎤=노트)

import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

SRC = "살리자_발표덱_v1.md"
OUT = "살리자_발표덱.pptx"

# ── 브랜드 토큰 ──
PRIMARY   = RGBColor(0x45, 0x70, 0xE8)
PRIM_LT   = RGBColor(0x93, 0xB2, 0xF8)
INK       = RGBColor(0x18, 0x18, 0x1E)
SUB       = RGBColor(0x42, 0x42, 0x4E)
MUTED     = RGBColor(0x80, 0x80, 0x8C)
BG        = RGBColor(0xF5, 0xF5, 0xFA)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
CODEBG    = RGBColor(0x12, 0x1A, 0x2E)
CODEINK   = RGBColor(0xCF, 0xDB, 0xFF)
CHIPBG    = RGBColor(0xEC, 0xF0, 0xFD)
FONT      = "Apple SD Gothic Neo"
MONO      = "Menlo"

EMU_W, EMU_H = Inches(13.333), Inches(7.5)


def strip_md(s):
    s = re.sub(r"\*\*(.+?)\*\*", r"\1", s)
    s = re.sub(r"\*(.+?)\*", r"\1", s)
    s = re.sub(r"`(.+?)`", r"\1", s)
    return s.strip()


def parse(md):
    slides = []
    cur = None
    in_code = False
    code = []
    for raw in md.splitlines():
        line = raw.rstrip()
        m = re.match(r"^##\s+S\d+\.\s+(.*)$", line)
        if m:
            if cur:
                slides.append(cur)
            cur = {"title": strip_md(m.group(1)), "items": [], "notes": []}
            continue
        if cur is None:
            continue
        if line.strip().startswith("```"):
            if in_code:
                cur["items"].append(("mono", "\n".join(code)))
                code = []
            in_code = not in_code
            continue
        if in_code:
            code.append(raw)
            continue
        if not line.strip():
            continue
        if "🎤" in line:  # 발표 멘트 → 슬라이드 노트 (- 🎤 / > 🎤 둘 다 허용)
            note = line.split("🎤", 1)[1].strip().strip('"').strip()
            cur["notes"].append(strip_md(note))
            continue
        mimg = re.match(r"^!\[(.*?)\]\((.+?)\)$", line.strip())
        if mimg:
            cur["items"].append(("image", (mimg.group(2).strip(), strip_md(mimg.group(1)))))
            continue
        if line.startswith(">"):
            cur["items"].append(("quote", strip_md(line[1:])))
            continue
        if line.startswith("- "):
            cur["items"].append(("bullet", strip_md(line[2:])))
            continue
        if re.match(r"^\*\*.+\*\*$", line.strip()):
            cur["items"].append(("lead", strip_md(line)))
            continue
        cur["items"].append(("para", strip_md(line)))
    if cur:
        slides.append(cur)
    return slides


def set_run(r, text, size, color, bold=False, font=FONT):
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.name = font


def add_rect(slide, x, y, w, h, color, line=None):
    sp = slide.shapes.add_shape(1, x, y, w, h)  # rectangle
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
    sp.shadow.inherit = False
    return sp


def textbox(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tb, tf


def add_notes(slide, notes):
    if not notes:
        return
    ns = slide.notes_slide
    ns.notes_text_frame.text = "\n".join(f"• {n}" for n in notes)


PHONE_AR = 390 / 844  # 폰 스크린샷 가로/세로 비


def add_phone_shots(slide, imgs):
    """우측 영역에 폰 스크린샷을 가로로 나란히 배치."""
    area_x = Inches(7.25)
    area_w = Inches(5.55)
    gap = Inches(0.18)
    n = len(imgs)
    # n개가 area_w 안에 들어가도록 높이 산정
    max_h = {1: 4.7, 2: 4.5, 3: 3.95}.get(n, 3.95)
    w_by_width = (area_w - gap * (n - 1)) / n
    h_by_width = w_by_width / PHONE_AR
    h = min(Emu(int(h_by_width)), Inches(max_h))
    w = Emu(int(h * PHONE_AR))
    total_w = w * n + gap * (n - 1)
    start_x = area_x + (area_w - total_w) / 2
    y = Inches(2.25)
    for i, (path, alt) in enumerate(imgs):
        x = start_x + (w + gap) * i
        if os.path.exists(path):
            slide.shapes.add_picture(path, x, y, height=h)
        # 캡션
        cap_y = y + h + Inches(0.05)
        tb, tf = textbox(slide, x - Inches(0.1), cap_y, w + Inches(0.2), Inches(0.5))
        cp = tf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        set_run(cp.add_run(), alt, 9, MUTED)


def build_title(prs, s):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, EMU_W, EMU_H, PRIMARY)
    add_rect(slide, 0, Inches(6.7), EMU_W, Inches(0.8), RGBColor(0x35, 0x5C, 0xD0))
    # 로고 마크
    mark = add_rect(slide, Inches(1.0), Inches(1.6), Inches(0.95), Inches(0.95), WHITE)
    mtf = mark.text_frame
    mtf.word_wrap = False
    p = mtf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_run(p.add_run(), "살", 34, PRIMARY, bold=True)
    mark.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    # 타이틀
    tb, tf = textbox(slide, Inches(1.0), Inches(2.9), Inches(11.3), Inches(2.6))
    p = tf.paragraphs[0]
    set_run(p.add_run(), "살리자 AI", 64, WHITE, bold=True)
    p2 = tf.add_paragraph()
    p2.space_before = Pt(6)
    set_run(p2.add_run(), "소상공인을 살리는 AI 파트너", 26, PRIM_LT, bold=True)
    p3 = tf.add_paragraph()
    p3.space_before = Pt(18)
    set_run(p3.add_run(), "데이터 인풋 → 분석 → 액션, 카카오에서 끝까지", 20, WHITE)
    # 하단 메타
    tb2, tf2 = textbox(slide, Inches(1.0), Inches(6.05), Inches(11.3), Inches(0.6))
    p = tf2.paragraphs[0]
    set_run(p.add_run(), "카카오 X SSAFY 해커톤 · 2026", 14, RGBColor(0xD8, 0xE2, 0xFF))
    add_notes(slide, s["notes"])


def build_content(prs, s, idx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, EMU_W, EMU_H, WHITE)
    add_rect(slide, 0, 0, Inches(0.16), EMU_H, PRIMARY)  # 좌측 액센트
    # 넘버 칩
    chip = add_rect(slide, Inches(0.6), Inches(0.5), Inches(1.0), Inches(0.42), CHIPBG)
    ctf = chip.text_frame
    ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
    cp = ctf.paragraphs[0]
    cp.alignment = PP_ALIGN.CENTER
    set_run(cp.add_run(), f"S{idx:02d}", 13, PRIMARY, bold=True)
    # 타이틀
    tb, tf = textbox(slide, Inches(0.6), Inches(1.05), Inches(12.1), Inches(1.1))
    p = tf.paragraphs[0]
    set_run(p.add_run(), s["title"], 30, INK, bold=True)
    # 이미지(폰 스크린샷) 있으면 우측 배치 + 본문 컬럼 축소
    imgs = [v for k, v in s["items"] if k == "image"]
    text_w = Inches(6.5) if imgs else Inches(12.1)
    # 본문
    tb2, tf2 = textbox(slide, Inches(0.62), Inches(2.15), text_w, Inches(4.7))
    first = True
    for kind, text in s["items"]:
        if kind in ("mono", "image"):
            continue  # 별도 영역
        p = tf2.paragraphs[0] if first else tf2.add_paragraph()
        first = False
        if kind == "bullet":
            p.space_after = Pt(7)
            r = p.add_run()
            set_run(r, "•  " + text, 17, SUB)
        elif kind == "lead":
            p.space_after = Pt(10)
            set_run(p.add_run(), text, 22, PRIMARY, bold=True)
        elif kind == "quote":
            p.space_after = Pt(10)
            set_run(p.add_run(), text, 18, INK, bold=True)
        else:  # para
            p.space_after = Pt(8)
            set_run(p.add_run(), text, 16, MUTED)
    # 코드/다이어그램 박스 (있으면 우측 정렬 큰 박스)
    monos = [t for k, t in s["items"] if k == "mono"]
    if monos:
        box = add_rect(slide, Inches(0.62), Inches(3.0), Inches(12.1), Inches(3.4), CODEBG)
        btf = box.text_frame
        btf.word_wrap = True
        btf.margin_left = Inches(0.25)
        btf.margin_top = Inches(0.18)
        lines = monos[0].split("\n")
        for i, ln in enumerate(lines):
            p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
            set_run(p.add_run(), ln if ln.strip() else " ", 14, CODEINK, font=MONO)
    # 폰 스크린샷 — 우측에 가로로 나란히
    if imgs:
        add_phone_shots(slide, imgs)
    # 푸터
    tb3, tf3 = textbox(slide, Inches(0.6), Inches(7.0), Inches(12.1), Inches(0.4))
    p = tf3.paragraphs[0]
    set_run(p.add_run(), "살리자 AI", 11, MUTED, bold=True)
    r = p.add_run()
    set_run(r, f"        ·        {idx:02d}", 11, MUTED)
    add_notes(slide, s["notes"])


def main():
    with open(SRC, encoding="utf-8") as f:
        slides = parse(f.read())
    prs = Presentation()
    prs.slide_width = EMU_W
    prs.slide_height = EMU_H
    for i, s in enumerate(slides, start=1):
        if i == 1:
            build_title(prs, s)
        else:
            build_content(prs, s, i)
    prs.save(OUT)
    print(f"OK · {len(slides)} slides → {OUT}")


if __name__ == "__main__":
    main()
